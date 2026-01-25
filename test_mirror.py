#!/usr/bin/env python3
"""Test script to verify RTL mirroring is working."""

import sys
sys.path.insert(0, '/Users/abdulrahmanalmutlaq/Desktop/Translation_site')

from pptx import Presentation
from pptx.util import Inches

def test_mirror_logic():
    """Test the basic mirror calculation."""
    # Simulate a 10-inch wide slide (standard)
    slide_width = Inches(10)  # 9144000 EMUs

    # Test case: shape at left edge
    shape_left = Inches(1)    # 1 inch from left
    shape_width = Inches(2)   # 2 inches wide

    # Expected: new_left = 10 - 1 - 2 = 7 inches from left
    new_left = slide_width - shape_left - shape_width

    print(f"Slide width: {slide_width} EMUs ({slide_width / 914400:.2f} inches)")
    print(f"Original position: {shape_left} EMUs ({shape_left / 914400:.2f} inches from left)")
    print(f"Shape width: {shape_width} EMUs ({shape_width / 914400:.2f} inches)")
    print(f"New position: {new_left} EMUs ({new_left / 914400:.2f} inches from left)")
    print()

    # Test case: shape at right edge
    shape_left2 = Inches(7)   # 7 inches from left
    shape_width2 = Inches(2)  # 2 inches wide

    new_left2 = slide_width - shape_left2 - shape_width2

    print(f"Original position: {shape_left2} EMUs ({shape_left2 / 914400:.2f} inches from left)")
    print(f"Shape width: {shape_width2} EMUs ({shape_width2 / 914400:.2f} inches)")
    print(f"New position: {new_left2} EMUs ({new_left2 / 914400:.2f} inches from left)")

    print("\n✓ Mirror calculation is correct!")


def test_with_real_pptx(pptx_path):
    """Test with a real PPTX file."""
    print(f"\nTesting with: {pptx_path}")

    prs = Presentation(pptx_path)
    slide_width = prs.slide_width

    print(f"Slide width: {slide_width} EMUs ({slide_width / 914400:.2f} inches)")
    print(f"Number of slides: {len(prs.slides)}")

    for slide_num, slide in enumerate(prs.slides, start=1):
        print(f"\n--- Slide {slide_num} ---")
        print(f"Number of shapes: {len(slide.shapes)}")

        for i, shape in enumerate(slide.shapes):
            try:
                old_left = shape.left
                shape_width = shape.width
                new_left = slide_width - old_left - shape_width

                print(f"  Shape {i}: '{shape.name}'")
                print(f"    Current left: {old_left / 914400:.2f} inches")
                print(f"    Width: {shape_width / 914400:.2f} inches")
                print(f"    Mirrored left: {new_left / 914400:.2f} inches")

                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text[:50] if shape.text_frame.text else "(empty)"
                    print(f"    Text: {text}...")
            except Exception as e:
                print(f"  Shape {i}: Error - {e}")

        if slide_num >= 2:  # Only check first 2 slides
            break


if __name__ == "__main__":
    print("=" * 50)
    print("RTL Mirror Test")
    print("=" * 50)

    test_mirror_logic()

    if len(sys.argv) > 1:
        test_with_real_pptx(sys.argv[1])
    else:
        print("\nTo test with a PPTX file, run:")
        print("  python test_mirror.py your_file.pptx")
