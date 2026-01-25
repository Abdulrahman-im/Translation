"""PPTX text extraction service."""

import re
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape
from typing import List, Tuple, Optional, Set


def parse_slide_range(range_str: str, max_slides: int) -> Set[int]:
    """
    Parse a slide range string into a set of slide numbers.

    Supports formats like:
    - "1-10" (range)
    - "1,3,5" (specific slides)
    - "1-5,8,10-12" (mixed)
    - "all" or empty string (all slides)

    Args:
        range_str: The range string to parse
        max_slides: Maximum number of slides in the presentation

    Returns:
        Set of slide numbers (1-indexed)
    """
    if not range_str or range_str.lower().strip() == "all":
        return set(range(1, max_slides + 1))

    slides = set()
    parts = range_str.replace(" ", "").split(",")

    for part in parts:
        if "-" in part:
            # Handle range like "1-10"
            match = re.match(r"(\d+)-(\d+)", part)
            if match:
                start = int(match.group(1))
                end = int(match.group(2))
                # Clamp to valid range
                start = max(1, min(start, max_slides))
                end = max(1, min(end, max_slides))
                slides.update(range(start, end + 1))
        else:
            # Handle single number
            try:
                num = int(part)
                if 1 <= num <= max_slides:
                    slides.add(num)
            except ValueError:
                pass

    return slides


def extract_text_from_shape(shape: BaseShape) -> List[str]:
    """Extract text from a single shape, handling different shape types."""
    texts = []

    # Handle grouped shapes recursively
    if isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            texts.extend(extract_text_from_shape(child_shape))
        return texts

    # Handle shapes with text frames
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text:
                    texts.append(text)

    # Handle tables
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            if text:
                                texts.append(text)

    return texts


def extract_text_from_pptx(
    file_path: str,
    slide_range: Optional[str] = None
) -> List[Tuple[int, str]]:
    """
    Extract all translatable text from a PowerPoint file.

    Args:
        file_path: Path to the PPTX file
        slide_range: Optional slide range (e.g., "1-10", "1,3,5", "1-5,8,10-12")

    Returns:
        List of (slide_number, text) tuples
    """
    prs = Presentation(file_path)
    total_slides = len(prs.slides)

    # Parse slide range
    slides_to_extract = parse_slide_range(slide_range or "", total_slides)

    results = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Skip slides not in the requested range
        if slide_num not in slides_to_extract:
            continue

        for shape in slide.shapes:
            texts = extract_text_from_shape(shape)
            for text in texts:
                # Skip empty or whitespace-only text
                if text and text.strip():
                    results.append((slide_num, text))

    return results


def get_slide_count(file_path: str) -> int:
    """Get the total number of slides in a PPTX file."""
    prs = Presentation(file_path)
    return len(prs.slides)
