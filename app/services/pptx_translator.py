"""
PPTX Translation Service.

Strategy:
1. Use PowerPoint (via AppleScript) for RTL mirroring - matches VBA exactly
2. Use python-pptx only for text translation

This ensures perfect RTL layout matching the VBA macro.
"""

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.oxml.ns import qn
from lxml import etree
from typing import Dict, Optional
import os
import platform

from .translator import translate_text
from .pptx_parser import parse_slide_range


def set_rtl_direction(slide):
    """Set RTL text direction for all shapes on a slide."""
    for shape in list(slide.shapes):
        if isinstance(shape, GroupShape):
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                try:
                    pPr = para._p.get_or_add_pPr()
                    pPr.set(qn('a:rtl'), '1')
                except:
                    pass

        elif shape.has_table:
            # Set table RTL
            try:
                tbl = shape.table._tbl
                tblPr = tbl.find(qn('a:tblPr'))
                if tblPr is None:
                    tblPr = etree.SubElement(tbl, qn('a:tblPr'))
                    tbl.insert(0, tblPr)
                tblPr.set('rtl', '1')
            except:
                pass

            # Set cell text RTL
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            try:
                                pPr = para._p.get_or_add_pPr()
                                pPr.set(qn('a:rtl'), '1')
                            except:
                                pass


def translate_slide_text(slide):
    """Translate all text in a slide (no layout changes)."""
    translations = []

    for shape in list(slide.shapes):
        if isinstance(shape, GroupShape):
            # After PowerPoint ungroups, there shouldn't be groups
            # But handle them just in case
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    orig = run.text.strip()
                    if orig:
                        trans = translate_text(orig)
                        run.text = trans
                        translations.append((orig, trans))

        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                orig = run.text.strip()
                                if orig:
                                    trans = translate_text(orig)
                                    run.text = trans
                                    translations.append((orig, trans))

    return translations


def translate_pptx_in_place(
    input_path: str,
    output_path: str,
    slide_range: Optional[str] = None,
    mirror_layout: bool = True
) -> Dict:
    """
    Translate PPTX with RTL mirroring.

    Process:
    1. If mirror_layout=True and on macOS: Use PowerPoint for mirroring (VBA logic)
    2. Then use python-pptx for translation only
    """
    print(f"\n{'='*60}")
    print("PPTX TRANSLATOR")
    print(f"{'='*60}")
    print(f"Input: {input_path}")
    print(f"Output: {output_path}")
    print(f"Mirror: {mirror_layout}")

    # Get slide info from original file
    prs_info = Presentation(input_path)
    total_slides = len(prs_info.slides)
    slides_to_process = parse_slide_range(slide_range or "", total_slides)
    print(f"Slides to process: {slides_to_process}")

    # STEP 1: Mirror using PowerPoint (if enabled and on macOS)
    if mirror_layout and platform.system() == 'Darwin':
        try:
            from .powerpoint_mirror import mirror_with_powerpoint, check_powerpoint_available

            if check_powerpoint_available():
                print(f"\n[STEP 1] Mirroring with PowerPoint...")
                mirror_with_powerpoint(input_path, output_path, list(slides_to_process))
                print(f"[STEP 1] PowerPoint mirroring complete!")
            else:
                print(f"\n[STEP 1] PowerPoint not available, copying file...")
                import shutil
                shutil.copy2(input_path, output_path)
        except Exception as e:
            print(f"[STEP 1] PowerPoint mirroring failed: {e}")
            print(f"[STEP 1] Falling back to copy only...")
            import shutil
            shutil.copy2(input_path, output_path)
    else:
        # No mirroring requested, just copy
        print(f"\n[STEP 1] Copying file (no mirroring)...")
        import shutil
        shutil.copy2(input_path, output_path)

    # STEP 2: Translate text using python-pptx
    print(f"\n[STEP 2] Translating text with python-pptx...")
    prs = Presentation(output_path)
    all_translations = []
    processed_slides = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_num not in slides_to_process:
            continue

        processed_slides += 1
        print(f"\n  Slide {slide_num}:")

        # Set RTL text direction (PowerPoint did position mirroring)
        if mirror_layout:
            set_rtl_direction(slide)
            print(f"    Set RTL text direction")

        slide_translations = translate_slide_text(slide)
        print(f"    Translated {len(slide_translations)} text items")

        for orig, trans in slide_translations:
            all_translations.append({
                "slide": slide_num,
                "original": orig,
                "translated": trans
            })

    # Save translated presentation
    print(f"\n[STEP 3] Saving...")
    prs.save(output_path)

    print(f"\n{'='*60}")
    print(f"Done! {len(all_translations)} items translated in {processed_slides} slides.")
    print(f"{'='*60}\n")

    return {
        "total_slides": total_slides,
        "processed_slides": processed_slides,
        "total_translations": len(all_translations),
        "translations": all_translations
    }


def translate_pptx_with_options(
    input_path: str,
    pptx_output_path: Optional[str] = None,
    excel_output_path: Optional[str] = None,
    slide_range: Optional[str] = None,
    mirror_layout: bool = True,
    output_excel: bool = False
) -> Dict:
    """Translate PPTX with flexible output options."""
    from .excel_writer import create_excel_file

    result = {
        "pptx_generated": False,
        "excel_generated": False,
        "pptx_path": None,
        "excel_path": None
    }

    if pptx_output_path:
        translation_result = translate_pptx_in_place(
            input_path, pptx_output_path, slide_range, mirror_layout
        )
        result["pptx_generated"] = True
        result["pptx_path"] = pptx_output_path
        result["total_slides"] = translation_result["total_slides"]
        result["processed_slides"] = translation_result["processed_slides"]
        result["total_translations"] = translation_result["total_translations"]
        result["translations"] = translation_result["translations"]

    if output_excel and excel_output_path and "translations" in result:
        excel_data = [
            (t["slide"], t["original"], t["translated"])
            for t in result["translations"]
        ]
        create_excel_file(excel_data, excel_output_path)
        result["excel_generated"] = True
        result["excel_path"] = excel_output_path

    return result
